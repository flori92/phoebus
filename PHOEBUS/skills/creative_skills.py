from PHOEBUS.skills.registry import skill
from pptx import Presentation
import os
import asyncio
from PHOEBUS.config import BASE_DIR

@skill(
    "create_presentation",
    risk="medium",
    help_text="Génère une présentation PowerPoint (.pptx) sur un sujet donné",
    describe=lambda d: f"Générer une présentation PowerPoint sur : {d.get('topic')}"
)
async def create_presentation(data: dict):
    topic = data.get("topic", "Sujet Inconnu")
    content = data.get("slides", []) # Liste de dict {"title": "...", "body": "..."}
    
    if not content:
        return "Je n'ai pas assez d'informations pour créer les diapositives."

    try:
        prs = Presentation()
        
        # Diapositive de titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = topic
        subtitle.text = "Généré par PHOEBUS 3.0"

        # Ajout des diapositives de contenu
        for item in content:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = item.get("title", "Sans titre")
            body_shape = slide.placeholders[1]
            body_shape.text = item.get("body", "")

        filename = f"Presentation_{topic.replace(' ', '_')[:20]}.pptx"
        save_path = os.path.join(os.path.expanduser("~/Downloads"), filename)
        
        prs.save(save_path)
        
        # Ouverture automatique
        from PHOEBUS.utils import open_uri
        open_uri(f"file://{save_path}")
        
        return f"La présentation sur '{topic}' a été générée et enregistrée dans vos Téléchargements."
    except Exception as e:
        return f"Erreur lors de la création du PowerPoint : {e}"

@skill(
    "generate_image",
    risk="low",
    help_text="Génère une image à partir d'une description textuelle",
    describe=lambda d: f"Générer une image pour : {d.get('prompt')}"
)
async def generate_image(data: dict):
    prompt = data.get("prompt")
    if not prompt:
        return "Quel type d'image voulez-vous que je génère ?"

    # Pour l'instant, on utilise une API gratuite et rapide (Pollinations ou similaire)
    # car les modèles diffusers locaux sont très lourds à charger au boot
    try:
        import requests
        import urllib.parse
        encoded_prompt = urllib.parse.quote(prompt)
        url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&nologo=true"
        
        save_path = os.path.join(os.path.expanduser("~/Downloads"), f"Image_{int(asyncio.get_event_loop().time())}.jpg")
        
        response = requests.get(url)
        if response.status_code == 200:
            with open(save_path, 'wb') as f:
                f.write(response.content)
            
            from PHOEBUS.utils import open_uri
            open_uri(f"file://{save_path}")
            return f"L'image pour '{prompt}' a été générée. Elle est ouverte sur votre écran."
        else:
            return "Le serveur de génération d'images est indisponible."
    except Exception as e:
        return f"Erreur lors de la génération d'image : {e}"
